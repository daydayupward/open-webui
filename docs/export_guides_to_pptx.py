#!/usr/bin/env python3
"""Export the two HTML guide decks to editable PowerPoint files."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
IMAGE_DIR = ROOT / "docs" / "image"
OUT_DIR = ROOT / "docs"

BG = RGBColor(7, 17, 31)
SURFACE = RGBColor(16, 38, 60)
INK = RGBColor(244, 248, 252)
MUTED = RGBColor(165, 184, 202)
CYAN = RGBColor(88, 217, 255)
BLUE = RGBColor(120, 147, 255)
GREEN = RGBColor(99, 230, 190)
AMBER = RGBColor(255, 193, 92)
PINK = RGBColor(236, 140, 255)
LINE = RGBColor(48, 75, 101)

W = Inches(13.333)
H = Inches(7.5)


def add_background(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_text(slide, text, x, y, w, h, size=24, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, eyebrow, title, subtitle=None):
    add_text(slide, eyebrow.upper(), 0.7, 0.45, 11.8, 0.3, 10, CYAN, True)
    add_text(slide, title, 0.7, 0.88, 11.9, 0.72, 28, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.65, 11.5, 0.5, 14, MUTED)


def add_footer(slide, label, number):
    add_text(slide, label.upper(), 0.7, 7.13, 5.5, 0.2, 8, RGBColor(109, 135, 159), True)
    add_text(slide, f"{number:02d}", 12.1, 7.13, 0.5, 0.2, 8, RGBColor(109, 135, 159), True, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, heading, body, accent=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    add_text(slide, heading, x + 0.22, y + 0.18, w - 0.44, 0.35, 15, accent, True)
    add_text(slide, body, x + 0.22, y + 0.62, w - 0.44, h - 0.78, 12, MUTED)
    return shape


def add_image(slide, filename, x, y, w, h):
    path = IMAGE_DIR / filename
    if not path.exists():
        add_text(slide, f"Missing image: {filename}", x, y, w, h, 14, AMBER)
        return
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def new_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    return prs, blank


def chiprag_deck():
    prs, blank = new_presentation()

    slide = prs.slides.add_slide(blank); add_background(slide)
    add_text(slide, "CHIP-RAG · ARCHITECTURE GUIDE", .7, .65, 9, .3, 11, CYAN, True)
    add_text(slide, "让芯片设计知识\n可检索、可验证、可追溯", .7, 1.25, 8.5, 1.8, 34, INK, True)
    add_text(slide, "面向 PDK 规则、EDA 脚本与设计指标的多专家 RAG 系统", .75, 3.35, 8.7, .55, 17, MUTED)
    add_text(slide, "核心主张：领域路由决定检索边界，Self-RAG 决定回答可信度，引用抽屉让证据可回溯。", .75, 4.35, 8.8, .75, 15, MUTED)
    add_footer(slide, "Architecture & Workflow", 1)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "01 · System architecture", "四层架构：从接入到知识库")
    add_image(slide, "arch_overview_cn.png", .7, 2.05, 7.25, 4.5)
    add_card(slide, 8.35, 2.05, 4.2, 1.0, "接入层", "Open WebUI · SSE", CYAN)
    add_card(slide, 8.35, 3.2, 4.2, 1.0, "编排层", "Supervisor · Expert · Finalizer", BLUE)
    add_card(slide, 8.35, 4.35, 4.2, 1.0, "检索层", "查询扩展 · 向量检索 · 重排序", GREEN)
    add_card(slide, 8.35, 5.5, 4.2, 1.0, "存储层", "PostgreSQL + pgvector", AMBER)
    add_footer(slide, "Architecture", 2)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "02 · Retrieval workflow", "RAG 工作流：先缩小范围，再丰富上下文")
    add_image(slide, "rag_workflow_cn.png", .75, 1.75, 11.8, 4.55)
    add_text(slide, "查询扩展  →  元数据硬过滤  →  向量召回 50  →  父文本注入  →  重排序 Top 10-15", .9, 6.45, 11.5, .35, 14, CYAN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Retrieval Pipeline", 3)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "03 · Grounding & context", "Parent-Child：精确匹配 + 丰富证据")
    add_image(slide, "chunking_strategy_cn.png", .7, 1.85, 7.2, 4.8)
    add_card(slide, 8.25, 1.9, 4.3, 1.15, "子块约 300 tokens", "适合向量匹配与精确定位。", CYAN)
    add_card(slide, 8.25, 3.25, 4.3, 1.15, "父块约 2000 tokens", "提供完整规则、表格和上下文。", BLUE)
    add_card(slide, 8.25, 4.6, 4.3, 1.15, "命中后注入 parent_text", "避免回答只看到短片段。", GREEN)
    add_footer(slide, "Chunking Strategy", 4)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "04 · Self-RAG & expert loop", "Self-RAG：回答不是终点，验证才是闭环")
    add_card(slide, .7, 2.0, 3.8, 3.45, "PDK Expert · 规则问答", "按工艺节点、类别与规则语义检索；幻觉或不完整时重写答案。\n\nDRC / LVS\n规则引用", CYAN)
    add_card(slide, 4.8, 2.0, 3.8, 3.45, "EDA Expert · 脚本生成", "Tcl/Skill 生成后经过程序化 Lint 与 LLM Review，再进入精炼循环。\n\n语法检查\n危险命令拦截", GREEN)
    add_card(slide, 8.9, 2.0, 3.8, 3.45, "Metrics Expert · 指标分析", "Text-to-SQL 只允许只读查询，并结合项目文档完成总结。\n\nSQL Guardrails\n单位规范", AMBER)
    add_text(slide, "验证器异常默认拒绝（False），并提供一次生成重试；长期可引入可观测评估集。", .9, 6.1, 11.5, .45, 14, MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Expert Loops", 5)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "05 · Positioning", "相对主流 RAG，Chip-RAG 的差异化")
    add_image(slide, "sota_comparison_cn.png", .75, 1.65, 11.8, 4.75)
    add_text(slide, "领域适配  ·  证据闭环  ·  工程安全", .9, 6.55, 11.5, .35, 15, CYAN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Differentiation", 6)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "06 · Roadmap", "下一步：从可用走向可度量")
    add_card(slide, .7, 1.9, 3.75, 3.7, "短期 P0", "混合检索（BM25 + 向量）\n对话摘要 + 窗口\n统一引用数据契约\n建立回归查询集", CYAN)
    add_card(slide, 4.8, 1.9, 3.75, 3.7, "中期 P1", "HyDE 查询扩展\n缓存层\n自动化评估流水线\n按任务风险动态选择验证强度", BLUE)
    add_card(slide, 8.9, 1.9, 3.75, 3.7, "长期方向", "多模态图文理解\n知识图谱\n工具调用\n跨项目知识治理", GREEN)
    add_text(slide, "建议指标：召回率 > 95% · 幻觉率 < 2% · 首 token 延迟 < 1s", .9, 6.2, 11.5, .4, 15, AMBER, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Roadmap & Measurement", 7)

    path = OUT_DIR / "Chip-RAG_Architecture_Guide.pptx"
    prs.save(path)
    return path


def workflows_deck():
    prs, blank = new_presentation()

    slide = prs.slides.add_slide(blank); add_background(slide)
    add_text(slide, "AGENT ENGINEERING · GUIDE", .7, .65, 9, .3, 11, CYAN, True)
    add_text(slide, "三大主流\nAI Agent 工作流", .7, 1.25, 8.5, 1.8, 36, INK, True)
    add_text(slide, "Matt Pocock Skills × Addy Osmani Agent Skills × obra Superpowers", .75, 3.35, 10, .55, 17, MUTED)
    add_text(slide, "从“能写代码”走向“可复用、可验证、可交付”。", .75, 4.35, 8.8, .5, 15, MUTED)
    add_footer(slide, "Workflow Comparison", 1)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "01 · Why workflow?", "模型能力之外，流程纪律决定交付质量")
    add_text(slide, "“优秀的 Agent 不是更会生成，而是更少跳过关键步骤。”", .8, 2.3, 7.1, 1.0, 25, INK, True)
    add_text(slide, "三套方案都把工程实践编码成可触发的技能，但在自由度、强制程度和多 Agent 协同上取舍不同。", .8, 3.65, 7.1, .95, 15, MUTED)
    for i, (h, b) in enumerate([("对齐", "理解目标与边界"), ("计划", "拆成可验证任务"), ("反馈", "测试、审查、修正"), ("交付", "确认结果再结束")]):
        add_card(slide, 8.45, 1.85 + i * 1.15, 4.1, .85, h, b, CYAN)
    add_footer(slide, "The Problem", 2)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "02 · Three philosophies", "三种工作流哲学")
    add_card(slide, .7, 2.0, 3.8, 3.7, "Matt Pocock Skills", "可组合的工程纪律\n\n领域建模 · 深模块 · TDD · 代码审查\n\n轻量 / 可组合 / 模型无关", AMBER)
    add_card(slide, 4.8, 2.0, 3.8, 3.7, "Addy Osmani Agent Skills", "完整生命周期\n\nDefine → Plan → Build → Test → Review → Ship\n\n验证门 / 跨平台 / 可复用", GREEN)
    add_card(slide, 8.9, 2.0, 3.8, 3.7, "obra Superpowers", "强制式工程方法论\n\nBrainstorming · Worktree · Plan · TDD · Subagent Review\n\n强约束 / 多 Agent / 可审计", PINK)
    add_footer(slide, "Three Approaches", 3)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "03 · Lifecycle comparison", "同一条主线，不同的控制面")
    add_image(slide, "agent_workflows_comparison_cn.png", .75, 1.65, 11.8, 4.95)
    add_footer(slide, "One Lifecycle · Three Control Styles", 4)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "04 · Comparison matrix", "能力矩阵：从“建议”到“门禁”")
    add_card(slide, .8, 2.0, 5.8, 3.8, "流程强度", "Matt Pocock：高灵活度 · 中等约束\n\nAddy Osmani：生命周期完整 · 验证明确\n\nSuperpowers：强制流程 · 协同最完整\n\n共同底线：计划、测试、审查、证据", CYAN)
    add_text(slide, "灵活性", 7.2, 2.0, 2.5, .3, 14, MUTED); add_text(slide, "高", 11.9, 2.0, .4, .3, 12, CYAN, True, align=PP_ALIGN.RIGHT)
    add_text(slide, "过程强制", 7.2, 3.0, 2.5, .3, 14, MUTED); add_text(slide, "中高", 11.5, 3.0, .8, .3, 12, CYAN, True, align=PP_ALIGN.RIGHT)
    add_text(slide, "多 Agent 编排", 7.2, 4.0, 2.5, .3, 14, MUTED); add_text(slide, "强", 11.9, 4.0, .4, .3, 12, CYAN, True, align=PP_ALIGN.RIGHT)
    for y, width in [(2.42, 4.6), (3.42, 3.9), (4.42, 4.2)]:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(y), Inches(width), Inches(.16)); s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()
    add_text(slide, "示意性相对强度，用于选择工作流，不代表官方 benchmark。", 7.2, 5.35, 5.0, .45, 11, MUTED)
    add_footer(slide, "Trade-offs", 5)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "05 · Fit for Chip-RAG", "对 Chip-RAG 的启发：组合，而非照搬")
    add_image(slide, "arch_overview_cn.png", .7, 2.0, 6.8, 4.25)
    add_card(slide, 7.8, 2.0, 4.8, 1.2, "Matt Pocock", "领域语言、架构边界、轻量设计讨论", AMBER)
    add_card(slide, 7.8, 3.45, 4.8, 1.2, "Addy Osmani", "将功能拆成 Spec → Plan → Build → Verify → Review", GREEN)
    add_card(slide, 7.8, 4.9, 4.8, 1.2, "Superpowers", "高风险改动、跨文件重构、多 Agent 实施", PINK)
    add_footer(slide, "Application to Chip-RAG", 6)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "06 · Recommended operating model", "推荐工作流：四道门")
    for i, label in enumerate(["ALIGN", "PLAN", "BUILD", "VERIFY"]):
        x = .9 + i * 3.05
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.1), Inches(1.55), Inches(1.55)); s.fill.solid(); s.fill.fore_color.rgb = RGBColor(16, 55, 80); s.line.color.rgb = CYAN
        add_text(slide, label, x, 2.66, 1.55, .3, 12, CYAN, True, align=PP_ALIGN.CENTER)
        if i < 3: add_text(slide, "→", x + 1.8, 2.65, .55, .35, 24, MUTED, True, align=PP_ALIGN.CENTER)
    add_card(slide, .8, 4.25, 2.8, 1.4, "对齐", "目标、非目标、成功标准", CYAN)
    add_card(slide, 3.85, 4.25, 2.8, 1.4, "计划", "文件、接口、依赖、风险", BLUE)
    add_card(slide, 6.9, 4.25, 2.8, 1.4, "实现", "垂直切片、TDD、独立 Agent", GREEN)
    add_card(slide, 9.95, 4.25, 2.8, 1.4, "验证", "测试、审查、运行时证据", AMBER)
    add_footer(slide, "Operating Model", 7)

    slide = prs.slides.add_slide(blank); add_background(slide); add_title(slide, "07 · Takeaways", "选择建议")
    add_text(slide, "3 → 1", .8, 2.0, 5.0, 1.0, 60, CYAN, True)
    add_text(slide, "吸收三套方案的长处，形成团队自己的最小可行工作流。", .85, 3.35, 5.6, .8, 18, MUTED)
    add_card(slide, 7.0, 1.95, 5.3, 1.0, "小改动", "Matt Pocock：轻量、直接", AMBER)
    add_card(slide, 7.0, 3.1, 5.3, 1.0, "常规功能", "Addy：完整、可复用", GREEN)
    add_card(slide, 7.0, 4.25, 5.3, 1.0, "高风险任务", "Superpowers：强制、可审计", PINK)
    add_text(slide, "最终原则：证据优先于感觉。", .9, 6.2, 11.5, .4, 17, CYAN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "End · Discuss", 8)

    path = OUT_DIR / "Agent_Workflows_Comparison_Guide.pptx"
    prs.save(path)
    return path


if __name__ == "__main__":
    print(chiprag_deck())
    print(workflows_deck())
